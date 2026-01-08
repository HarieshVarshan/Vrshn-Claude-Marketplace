"""
Document Text Extraction Module
Supports multiple document formats: PDF, DOCX, XLSX, PPTX, ODT, TXT, MD, HTML, CSV
"""

import csv
import io
from pathlib import Path
from typing import Optional

# Supported file extensions
SUPPORTED_EXTENSIONS = {
    '.pdf': 'PDF Document',
    '.docx': 'Microsoft Word',
    '.doc': 'Microsoft Word (Legacy)',
    '.xlsx': 'Microsoft Excel',
    '.xls': 'Microsoft Excel (Legacy)',
    '.pptx': 'Microsoft PowerPoint',
    '.odt': 'OpenDocument Text',
    '.ods': 'OpenDocument Spreadsheet',
    '.odp': 'OpenDocument Presentation',
    '.txt': 'Plain Text',
    '.md': 'Markdown',
    '.html': 'HTML',
    '.htm': 'HTML',
    '.csv': 'CSV',
    '.json': 'JSON',
    '.xml': 'XML',
    '.rtf': 'Rich Text Format',
}


def get_supported_extensions() -> list[str]:
    """Return list of supported file extensions."""
    return list(SUPPORTED_EXTENSIONS.keys())


def is_supported(file_path: str) -> bool:
    """Check if file type is supported."""
    ext = Path(file_path).suffix.lower()
    return ext in SUPPORTED_EXTENSIONS


def get_file_type(file_path: str) -> str:
    """Get human-readable file type."""
    ext = Path(file_path).suffix.lower()
    return SUPPORTED_EXTENSIONS.get(ext, 'Unknown')


def extract_text(file_path: str) -> str:
    """
    Extract text from a document file.

    Args:
        file_path: Path to the document file

    Returns:
        Extracted text as a single string

    Raises:
        ValueError: If file type is not supported
        FileNotFoundError: If file doesn't exist
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower()

    extractors = {
        '.pdf': _extract_pdf,
        '.docx': _extract_docx,
        '.doc': _extract_doc_legacy,
        '.xlsx': _extract_xlsx,
        '.xls': _extract_xlsx,
        '.pptx': _extract_pptx,
        '.odt': _extract_odt,
        '.ods': _extract_ods,
        '.odp': _extract_odp,
        '.txt': _extract_text,
        '.md': _extract_markdown,
        '.html': _extract_html,
        '.htm': _extract_html,
        '.csv': _extract_csv,
        '.json': _extract_text,
        '.xml': _extract_xml,
        '.rtf': _extract_rtf,
    }

    extractor = extractors.get(ext)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {ext}")

    return extractor(file_path)


def extract_with_metadata(file_path: str) -> dict:
    """
    Extract text and metadata from a document file.

    Args:
        file_path: Path to the document file

    Returns:
        Dictionary with text, file_type, and metadata
    """
    path = Path(file_path)
    text = extract_text(file_path)

    return {
        "file_path": str(path.absolute()),
        "file_name": path.name,
        "file_type": get_file_type(file_path),
        "extension": path.suffix.lower(),
        "size_bytes": path.stat().st_size,
        "full_text": text,
        "char_count": len(text),
    }


# --- Individual Extractors ---

def _extract_pdf(file_path: str) -> str:
    """Extract text from PDF using PyMuPDF."""
    import fitz

    doc = fitz.open(file_path)
    text = ""
    for page_num, page in enumerate(doc):
        page_text = page.get_text()
        if page_text.strip():
            text += f"\n--- Page {page_num + 1} ---\n"
            text += page_text
    doc.close()
    return text


def _extract_docx(file_path: str) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document

    doc = Document(file_path)
    paragraphs = []

    # Extract paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)

    # Extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    return '\n\n'.join(paragraphs)


def _extract_doc_legacy(file_path: str) -> str:
    """Extract text from legacy DOC files."""
    # Try using antiword or catdoc if available, otherwise return error message
    import subprocess

    # Try antiword first
    try:
        result = subprocess.run(
            ['antiword', file_path],
            capture_output=True,
            text=True,
            timeout=30
        )
        if result.returncode == 0:
            return result.stdout
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    # Try catdoc
    try:
        result = subprocess.run(
            ['catdoc', file_path],
            capture_output=True,
            text=True,
            timeout=30
        )
        if result.returncode == 0:
            return result.stdout
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass

    return f"[Legacy .doc format - install 'antiword' or 'catdoc' to extract: {Path(file_path).name}]"


def _extract_xlsx(file_path: str) -> str:
    """Extract text from XLSX using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True)
    text_parts = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        text_parts.append(f"\n--- Sheet: {sheet_name} ---\n")

        rows = []
        for row in sheet.iter_rows():
            row_values = []
            for cell in row:
                if cell.value is not None:
                    row_values.append(str(cell.value))
            if row_values:
                rows.append(' | '.join(row_values))

        text_parts.append('\n'.join(rows))

    wb.close()
    return '\n'.join(text_parts)


def _extract_pptx(file_path: str) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation

    prs = Presentation(file_path)
    text_parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        text_parts.append(f"\n--- Slide {slide_num} ---\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)

            # Extract table text
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        slide_text.append(row_text)

        text_parts.append('\n'.join(slide_text))

    return '\n'.join(text_parts)


def _extract_odt(file_path: str) -> str:
    """Extract text from ODT using odfpy."""
    from odf import text as odf_text
    from odf.opendocument import load

    doc = load(file_path)
    paragraphs = []

    for para in doc.getElementsByType(odf_text.P):
        text = ""
        for node in para.childNodes:
            if hasattr(node, 'data'):
                text += node.data
            elif node.tagName == 'text:s':
                text += ' '
        if text.strip():
            paragraphs.append(text)

    return '\n\n'.join(paragraphs)


def _extract_ods(file_path: str) -> str:
    """Extract text from ODS using odfpy."""
    from odf.opendocument import load
    from odf import table as odf_table
    from odf import text as odf_text

    doc = load(file_path)
    text_parts = []

    for sheet in doc.getElementsByType(odf_table.Table):
        sheet_name = sheet.getAttribute('name')
        text_parts.append(f"\n--- Sheet: {sheet_name} ---\n")

        rows = []
        for row in sheet.getElementsByType(odf_table.TableRow):
            row_values = []
            for cell in row.getElementsByType(odf_table.TableCell):
                cell_text = ""
                for para in cell.getElementsByType(odf_text.P):
                    for node in para.childNodes:
                        if hasattr(node, 'data'):
                            cell_text += node.data
                if cell_text.strip():
                    row_values.append(cell_text.strip())
            if row_values:
                rows.append(' | '.join(row_values))

        text_parts.append('\n'.join(rows))

    return '\n'.join(text_parts)


def _extract_odp(file_path: str) -> str:
    """Extract text from ODP using odfpy."""
    from odf.opendocument import load
    from odf import draw as odf_draw
    from odf import text as odf_text

    doc = load(file_path)
    text_parts = []
    slide_num = 0

    for page in doc.getElementsByType(odf_draw.Page):
        slide_num += 1
        text_parts.append(f"\n--- Slide {slide_num} ---\n")

        slide_text = []
        for frame in page.getElementsByType(odf_draw.Frame):
            for textbox in frame.getElementsByType(odf_draw.TextBox):
                for para in textbox.getElementsByType(odf_text.P):
                    para_text = ""
                    for node in para.childNodes:
                        if hasattr(node, 'data'):
                            para_text += node.data
                    if para_text.strip():
                        slide_text.append(para_text)

        text_parts.append('\n'.join(slide_text))

    return '\n'.join(text_parts)


def _extract_text(file_path: str) -> str:
    """Extract text from plain text files."""
    encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']

    for encoding in encodings:
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            continue

    # Fallback: read as binary and decode with errors='replace'
    with open(file_path, 'rb') as f:
        return f.read().decode('utf-8', errors='replace')


def _extract_markdown(file_path: str) -> str:
    """Extract text from Markdown, preserving structure."""
    return _extract_text(file_path)


def _extract_html(file_path: str) -> str:
    """Extract text from HTML using BeautifulSoup."""
    from bs4 import BeautifulSoup

    html_content = _extract_text(file_path)
    soup = BeautifulSoup(html_content, 'lxml')

    # Remove script and style elements
    for script in soup(["script", "style"]):
        script.decompose()

    # Get text
    text = soup.get_text(separator='\n')

    # Clean up whitespace
    lines = (line.strip() for line in text.splitlines())
    text = '\n'.join(line for line in lines if line)

    return text


def _extract_csv(file_path: str) -> str:
    """Extract text from CSV files."""
    text_content = _extract_text(file_path)

    try:
        reader = csv.reader(io.StringIO(text_content))
        rows = []
        for row in reader:
            rows.append(' | '.join(str(cell) for cell in row if cell))
        return '\n'.join(rows)
    except csv.Error:
        return text_content


def _extract_xml(file_path: str) -> str:
    """Extract text from XML files."""
    from bs4 import BeautifulSoup

    xml_content = _extract_text(file_path)
    soup = BeautifulSoup(xml_content, 'lxml-xml')

    # Get all text content
    text = soup.get_text(separator='\n')

    # Clean up whitespace
    lines = (line.strip() for line in text.splitlines())
    text = '\n'.join(line for line in lines if line)

    return text


def _extract_rtf(file_path: str) -> str:
    """Extract text from RTF files (basic extraction)."""
    import re

    content = _extract_text(file_path)

    # Basic RTF tag removal
    # Remove RTF control words
    text = re.sub(r'\\[a-z]+\d* ?', '', content)
    # Remove groups
    text = re.sub(r'[{}]', '', text)
    # Clean up whitespace
    text = re.sub(r'\s+', ' ', text)

    return text.strip()


def extract_from_multiple(file_paths: list[str]) -> dict[str, str]:
    """
    Extract text from multiple document files.

    Args:
        file_paths: List of paths to document files

    Returns:
        Dictionary mapping file paths to extracted text
    """
    results = {}
    for path in file_paths:
        try:
            results[path] = extract_text(path)
        except Exception as e:
            print(f"Error extracting {path}: {e}")
            results[path] = ""
    return results


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python document_extractor.py <file>")
        print(f"\nSupported formats: {', '.join(SUPPORTED_EXTENSIONS.keys())}")
        sys.exit(1)

    file_path = sys.argv[1]

    if not is_supported(file_path):
        print(f"Unsupported file type: {Path(file_path).suffix}")
        print(f"Supported: {', '.join(SUPPORTED_EXTENSIONS.keys())}")
        sys.exit(1)

    try:
        result = extract_with_metadata(file_path)
        print(f"File: {result['file_name']}")
        print(f"Type: {result['file_type']}")
        print(f"Size: {result['size_bytes']} bytes")
        print(f"Characters: {result['char_count']}")
        print("-" * 50)
        text = result['full_text']
        print(text[:2000] + "..." if len(text) > 2000 else text)
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)
